from pptx import Presentation

answer_time = [15,30,30,45,15]
sourcepath = r"C:\Users\Léo-Paul\Desktop\Sourcetest.pptx"
prs = Presentation(str(sourcepath))

list = []
ranges = 5 # nombres de temps différents prévus
for question_number in range (5):
    time = answer_time[question_number]
    range_of_this_time = (time // 15) - 1
    ranges_to_delete = [0,1,2]
    del ranges_to_delete[range_of_this_time]
    for k in ranges_to_delete :
        slide_to_delete = 1 + 3*question_number + k
        list.append(slide_to_delete)
        
list = sorted(list)
slides_to_delete = []

for i in range(len(list)-1,-1,-1):
    slides_to_delete.append(list[i])
        
for slide_to_delete in slides_to_delete:        
    rId = prs.slides._sldIdLst[slide_to_delete+1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_to_delete+1]
   
prs.save(r"C:\Users\Léo-Paul\Desktop\Testdel.pptx")