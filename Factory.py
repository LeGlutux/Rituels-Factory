from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from xlrd import open_workbook
import os


sourcepath = r"C:\Users\Léo-Paul\Desktop\Rituels Factory\Sources"


for t in range(4): # choix du niveau 0 = 6ème ... 3 = 3ème

    if t == 0:
        classe = '6ème'
    elif t == 1:
        classe = '5ème'
    elif t == 2:
        classe = '4ème'
    elif t == 3:
        classe = '3ème'
    
    pathout = r"C:\Users\Léo-Paul\Desktop\Rituels Factory" + "\-" + str(classe)
# Extraction du fichier excel data est la liste des listes lignes
# sheet.nrows = nombre de lignes

    book = open_workbook(str(sourcepath)+"\Rituels(Calculs).xlsx")
    sheet = book.sheet_by_index(t)

    column_index = 0
    data=[]
    for row_index in range(1, sheet.nrows):
        text = sheet.row_values(row_index, column_index, end_colx=6)
        data.append(text)

# Création des dossiers contenants les pptx
    
    FolderNumber = int((sheet.nrows - 1) / 5) # nombre de dossiers pour une classe donnée
    for i in range(FolderNumber):
        path = str(pathout)+"\-" +str(i+1)+ "- " + str(data[i*5][0])
        if not os.path.exists(path):
            os.mkdir(path)
        
# Position flèche slide 1
    
    Top_Fleche_1 = 1700808
    Top_Fleche_2 = 1196752
    Dist_Fleche = 969000
    

    for j in range(FolderNumber):
        for i in range(5):
            prs = Presentation(str(sourcepath)+'\Matrix ' + str(classe) + '.pptx')
            prs.slides[0].shapes[2].top = 1700808 + i*Dist_Fleche
            prs.slides[1].shapes[5].top = 1196752 + i*Dist_Fleche
            if i == 4:
# Explosion éval
                    
                shapes = prs.slides[0].shapes                   
                left = Inches(0.8)               
                top = Inches(5.5)               
                width = Inches(4.8)               
                height = Inches(1.5)               
                shape = shapes.add_shape(MSO_SHAPE.EXPLOSION1, left, top, width, height)               
                cursor_sp = shapes[0]._element               
                cursor_sp.addprevious(shapes[4]._element)                
                fill = shape.fill               
                fill.solid()               
                fill.fore_color.rgb = RGBColor(255, 255, 0)              
                line = shape.line               
                line.color.rgb = RGBColor(0, 0, 0)                
                line.color.brightness = 1.0  # % light                
                line.width = Pt(0.0)
                
# Suppression des slides en trop (temps de réponse)
  
                answer_time = sheet.row_values(j*5 + i + 1, 6, end_colx=11)
                answer_time.append(text)
                list = []
                ranges = 5 # nombres de temps différents prévus
                for question_number in range (5):
                    time = answer_time[question_number]
                    if time == 15 :
                        range_of_this_time = 0
                    if time == 15 :
                        range_of_this_time = 1
                    if time == 15 :
                        range_of_this_time = 2
                    if time == 15 :
                        range_of_this_time = 3
                    if time == 15 :
                        range_of_this_time = 4
                    else:
                        print(error)
                    ranges_to_delete = [0,1,2,3,4]
                    del ranges_to_delete[range_of_this_time]
                    for k in ranges_to_delete :
                        slide_to_delete = 1 + 5*question_number + k
                        list.append(slide_to_delete)
        
            list = sorted(list)
            slides_to_delete = []

            for i in range(len(list)-1,-1,-1):
                slides_to_delete.append(list[i])
        
            for slide_to_delete in slides_to_delete:        
                rId = prs.slides._sldIdLst[slide_to_delete+1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_to_delete+1]
   
    
# Slide de calcul
            for k in range(5):
                slide = prs.slides[2+k]
                image_path = r"C:\Users\Léo-Paul\Desktop\Rituels Factory\Sources\Images\\"+ str(classe) + "\\" + str(j+1) + "-"+str(i+1)+"-"+str(k+1)+".png"
                if os.path.exists(image_path): # Ajout d'images
                    shape = slide.shapes[0]
                    pic = shape.insert_picture(image_path)
             
                    
                text_frame = slide.shapes[7].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(data[5*j+i][k+1])
                font = run.font
                font.name = 'Calibri'
                if  len(data[5*j + i][k+1]) < 20 :
                    
                    font.size = Pt(68)
                    font.bold = True
                    
                else:
                    font.size = Pt(38)
                    font.bold = False
                    
                font.italic = None  # cause value to be inherited from theme
                font.color.theme_color = MSO_THEME_COLOR.DARK_1
                text_frame.margin_bottom = Inches(0.08)
                text_frame.margin_left = 0
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                text_frame.word_wrap = False
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                if i == 4 :
                    prs.save(str(pathout)+"\-" +str(j+1)+ "- " + str(data[j*5][0])+"\zÉvaluation.pptx")
                else:
                    prs.save(str(pathout)+"\-" +str(j+1)+ "- " + str(data[j*5][0])+"\Rituel "+str(i+1)+'.pptx')
                
                
